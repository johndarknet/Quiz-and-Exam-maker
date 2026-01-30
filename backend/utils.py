import random
from pptx import Presentation
import PyPDF2
import nltk
from nltk import word_tokenize, pos_tag, sent_tokenize

# download needed NLTK data (first run only)
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

try:
    nltk.data.find('taggers/averaged_perceptron_tagger')
except LookupError:
    nltk.download('averaged_perceptron_tagger')


def extract_text_from_pdf(file_path_or_stream):
    reader = PyPDF2.PdfReader(file_path_or_stream)
    texts = []
    for page in reader.pages:
        text = page.extract_text() or ""
        texts.append(text)
    return "\n".join(texts)


def extract_text_from_pptx(file_path_or_stream):
    prs = Presentation(file_path_or_stream)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text:
                    texts.append(shape.text)
    return "\n".join(texts)


def get_nouns(sentence):
    words = word_tokenize(sentence)
    tags = pos_tag(words)
    nouns = [w for w, t in tags if t.startswith("NN") and w.isalpha() and len(w) > 2]
    return nouns


def make_multiple_choice(sentence, all_nouns, qid):
    nouns = get_nouns(sentence)
    if not nouns:
        return None

    answer = random.choice(nouns)

    # distractors
    distractors = [n for n in all_nouns if n.lower() != answer.lower()]
    random.shuffle(distractors)
    choices = [answer] + distractors[:3]

    while len(choices) < 4:
        choices.append("Option" + str(len(choices)))

    random.shuffle(choices)

    return {
        "id": qid,
        "type": "multiple_choice",
        "question": f"In the sentence: \"{sentence}\" what word best completes the blank?\n{sentence.replace(answer, '_____', 1)}",
        "choices": choices,
        "answer": answer
    }


def make_fill_blank(sentence, qid):
    nouns = get_nouns(sentence)
    if not nouns:
        return None

    answer = random.choice(nouns)
    prompt = sentence.replace(answer, "_____", 1)

    return {
        "id": qid,
        "type": "fill_blank",
        "question": f"Fill in the blank:\n{prompt}",
        "answer": answer
    }


def make_enumeration(sentence, qid):
    # try to find comma separated list
    parts = [p.strip() for p in sentence.split(",")]
    if len(parts) < 3:
        return None

    items = parts[:3]
    question = "Enumerate the following items mentioned:\n" + ", ".join(items)

    return {
        "id": qid,
        "type": "enumeration",
        "question": question,
        "answer": items
    }


def simple_generate_quiz(text, max_questions=10):
    sentences = [s.strip() for s in sent_tokenize(text) if len(s.strip()) > 30]
    random.shuffle(sentences)

    # collect noun pool for distractors
    all_nouns = []
    for s in sentences:
        all_nouns.extend(get_nouns(s))

    questions = []
    qid = 1

    for s in sentences:
        if qid > max_questions:
            break

        qtype = random.choice(["mcq", "fill", "enum"])

        if qtype == "mcq":
            q = make_multiple_choice(s, all_nouns, qid)
        elif qtype == "fill":
            q = make_fill_blank(s, qid)
        else:
            q = make_enumeration(s, qid)

        if q is not None:
            questions.append(q)
            qid += 1

    return questions
